import logging
import shutil
import subprocess
import tempfile
from pathlib import Path

from langchain_community.document_loaders import PyPDFLoader, TextLoader
from langchain_core.documents import Document
from docx import Document as DocxDocument
from docx.oxml.ns import qn
from openpyxl import load_workbook
from pptx import Presentation

logger = logging.getLogger(__name__)

try:
    from office_oxide import Document as OfficeDocument
except ImportError:  # pragma: no cover - optional on unsupported Python builds
    OfficeDocument = None


def _office_page_metadata(**extra: object) -> dict:
    meta = {"parse_source": extra.pop("parse_source", "python_office")}
    if "page_number" in extra:
        meta["page_number"] = extra.pop("page_number")
    meta.update(extra)
    return meta


def _paragraph_has_page_break(paragraph) -> bool:
    for run in paragraph.runs:
        for child in run._element:
            if child.tag == qn("w:br") and child.get(qn("w:type")) == "page":
                return True
    return False


def _load_docx_with_pages(file_path: str) -> list[Document]:
    doc = DocxDocument(file_path)
    pages: dict[int, list[str]] = {1: []}
    current_page = 1

    for paragraph in doc.paragraphs:
        if _paragraph_has_page_break(paragraph):
            current_page += 1
            pages.setdefault(current_page, [])
        text = paragraph.text.strip()
        if text:
            pages.setdefault(current_page, []).append(text)

    docs: list[Document] = []
    for page_number, parts in sorted(pages.items()):
        if parts:
            docs.append(
                Document(
                    page_content="\n".join(parts),
                    metadata=_office_page_metadata(parse_source="python_docx", page_number=page_number),
                )
            )
    return docs


def _load_xlsx(file_path: str) -> list[Document]:
    workbook = load_workbook(file_path, read_only=True, data_only=True)
    docs: list[Document] = []
    for sheet_idx, sheet in enumerate(workbook.worksheets, start=1):
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell).strip() if cell is not None else "" for cell in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            docs.append(
                Document(
                    page_content="\n".join(rows),
                    metadata=_office_page_metadata(
                        parse_source="python_xlsx",
                        page_number=sheet_idx,
                        sheet_name=sheet.title,
                    ),
                )
            )
    workbook.close()
    return docs


def _load_xls_xlrd(file_path: str) -> list[Document]:
    import xlrd

    workbook = xlrd.open_workbook(file_path)
    docs: list[Document] = []
    for sheet_idx in range(workbook.nsheets):
        sheet = workbook.sheet_by_index(sheet_idx)
        rows: list[str] = []
        for row_idx in range(sheet.nrows):
            cells = [str(sheet.cell_value(row_idx, col_idx)).strip() for col_idx in range(sheet.ncols)]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            docs.append(
                Document(
                    page_content="\n".join(rows),
                    metadata=_office_page_metadata(
                        parse_source="python_xls",
                        page_number=sheet_idx + 1,
                        sheet_name=sheet.name,
                    ),
                )
            )
    return docs


def _load_pptx(file_path: str) -> list[Document]:
    presentation = Presentation(file_path)
    docs: list[Document] = []
    for slide_idx, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    parts.append(text)
        if parts:
            docs.append(
                Document(
                    page_content="\n".join(parts),
                    metadata=_office_page_metadata(
                        parse_source="python_pptx",
                        page_number=slide_idx,
                        slide_index=slide_idx,
                    ),
                )
            )
    return docs


def _ir_sheet_rows(sheet: dict) -> list[str]:
    rows: list[str] = []
    for row in sheet.get("rows") or []:
        if not isinstance(row, dict):
            continue
        cells = row.get("cells") or []
        values = [str(c.get("value", "")).strip() for c in cells if isinstance(c, dict)]
        if any(values):
            rows.append(" | ".join(values))
    return rows


def _ir_slide_text(slide: dict) -> str:
    parts: list[str] = []
    for block in slide.get("blocks") or []:
        if not isinstance(block, dict):
            continue
        text = str(block.get("text") or "").strip()
        if text:
            parts.append(text)
    return "\n".join(parts)


def _load_with_office_oxide(file_path: str, file_type: str) -> list[Document]:
    assert OfficeDocument is not None
    with OfficeDocument.open(file_path) as doc:
        ir = doc.to_ir()
        docs: list[Document] = []

        if file_type == "xls":
            sheets = ir.get("sheets") if isinstance(ir, dict) else None
            if isinstance(sheets, list) and sheets:
                for sheet_idx, sheet in enumerate(sheets, start=1):
                    if not isinstance(sheet, dict):
                        continue
                    rows = _ir_sheet_rows(sheet)
                    if rows:
                        docs.append(
                            Document(
                                page_content="\n".join(rows),
                                metadata=_office_page_metadata(
                                    parse_source="python_xls",
                                    page_number=sheet_idx,
                                    sheet_name=str(sheet.get("name") or f"Sheet{sheet_idx}"),
                                ),
                            )
                        )
            if docs:
                return docs

        if file_type == "ppt":
            slides = ir.get("slides") if isinstance(ir, dict) else None
            if isinstance(slides, list) and slides:
                for slide_idx, slide in enumerate(slides, start=1):
                    if not isinstance(slide, dict):
                        continue
                    text = _ir_slide_text(slide)
                    if text:
                        docs.append(
                            Document(
                                page_content=text,
                                metadata=_office_page_metadata(
                                    parse_source="python_ppt",
                                    page_number=slide_idx,
                                    slide_index=slide_idx,
                                ),
                            )
                        )
            if docs:
                return docs

        text = doc.plain_text().strip()
        if not text:
            return []
        return [
            Document(
                page_content=text,
                metadata=_office_page_metadata(parse_source=f"python_{file_type}", page_number=1),
            )
        ]


def _libreoffice_convert(file_path: str, output_ext: str, out_dir: str) -> str:
    if shutil.which("soffice") is None:
        raise RuntimeError("LibreOffice (soffice) is required for legacy Office parsing")
    cmd = [
        "soffice", "--headless", "--convert-to", output_ext,
        "--outdir", out_dir, file_path,
    ]
    subprocess.run(cmd, check=True, capture_output=True)
    converted = list(Path(out_dir).glob(f"*.{output_ext}"))
    if not converted:
        raise RuntimeError(f"LibreOffice conversion to {output_ext} produced no output")
    return str(converted[0])


def _load_legacy_office(file_path: str, file_type: str) -> list[Document]:
    if OfficeDocument is not None:
        try:
            return _load_with_office_oxide(file_path, file_type)
        except Exception:
            logger.exception("office-oxide failed for %s, trying fallback", file_path)

    if file_type == "xls":
        return _load_xls_xlrd(file_path)

    target_ext = {"doc": "docx", "ppt": "pptx"}[file_type]
    with tempfile.TemporaryDirectory(prefix="rag_office_") as temp_dir:
        converted_path = _libreoffice_convert(file_path, target_ext, temp_dir)
        if target_ext == "docx":
            return _load_docx_with_pages(converted_path)
        return _load_pptx(converted_path)


def load_with_python(file_path: str, file_type: str) -> list[Document]:
    if file_type == "pdf":
        loader = PyPDFLoader(file_path)
        docs = loader.load()
        for doc in docs:
            page = doc.metadata.get("page")
            if page is not None:
                doc.metadata["page_number"] = int(page) + 1
            doc.metadata.setdefault("parse_source", "python_pdf")
        return docs
    if file_type in {"txt", "md"}:
        loader = TextLoader(file_path, encoding="utf-8")
        return loader.load()
    if file_type == "docx":
        return _load_docx_with_pages(file_path)
    if file_type == "xlsx":
        return _load_xlsx(file_path)
    if file_type == "pptx":
        return _load_pptx(file_path)
    if file_type in {"doc", "xls", "ppt"}:
        return _load_legacy_office(file_path, file_type)
    raise ValueError(f"Python loader does not support file type: {file_type}")


def markdown_to_documents(markdown: str, *, source: str = "mineru") -> list[Document]:
    text = markdown.strip()
    if not text:
        return []
    return [Document(page_content=text, metadata={"parse_source": source})]
