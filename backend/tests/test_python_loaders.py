from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from openpyxl import Workbook
from pptx import Presentation

from app.parsing.python_loaders import _load_docx_with_pages, _load_legacy_office, load_with_python


def test_load_xlsx_has_page_number(tmp_path: Path):
    path = tmp_path / "data.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Sheet1"
    sheet["A1"] = "Name"
    sheet["B1"] = "Value"
    sheet["A2"] = "foo"
    sheet["B2"] = "bar"
    workbook.save(path)

    docs = load_with_python(str(path), "xlsx")
    assert len(docs) >= 1
    assert docs[0].metadata.get("page_number") == 1
    assert docs[0].metadata.get("sheet_name") == "Sheet1"


def test_load_pptx_has_page_number(tmp_path: Path):
    path = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Title slide"
    presentation.save(path)

    docs = load_with_python(str(path), "pptx")
    assert len(docs) == 1
    assert docs[0].metadata.get("page_number") == 1
    assert docs[0].metadata.get("slide_index") == 1


def test_load_docx_with_pages_splits_on_page_break(tmp_path: Path, monkeypatch):
    from docx import Document as DocxDocument
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    path = tmp_path / "doc.docx"
    doc = DocxDocument()
    doc.add_paragraph("Page one")
    para = doc.add_paragraph("Page two")
    run = para.add_run()
    br = OxmlElement("w:br")
    br.set(qn("w:type"), "page")
    run._r.append(br)
    run.add_text("After break")
    doc.save(path)

    docs = _load_docx_with_pages(str(path))
    assert len(docs) >= 1
    assert docs[0].metadata.get("page_number") == 1


def test_load_legacy_doc_uses_office_oxide(tmp_path: Path):
    path = tmp_path / "legacy.doc"
    path.write_bytes(b"fake")

    mock_doc = MagicMock()
    mock_doc.plain_text.return_value = "legacy content"
    mock_doc.to_ir.return_value = {}

    with patch("app.parsing.python_loaders.OfficeDocument") as mock_cls:
        mock_cls.open.return_value.__enter__.return_value = mock_doc
        docs = _load_legacy_office(str(path), "doc")

    assert len(docs) == 1
    assert docs[0].metadata.get("page_number") == 1
    assert "legacy content" in docs[0].page_content


def test_load_legacy_xls_uses_xlrd_when_no_office_oxide(tmp_path: Path, monkeypatch):
    path = tmp_path / "legacy.xls"
    path.write_bytes(b"fake")
    monkeypatch.setattr("app.parsing.python_loaders.OfficeDocument", None)

    with patch("app.parsing.python_loaders._load_xls_xlrd") as mock_xlrd:
        mock_xlrd.return_value = [
            __import__("langchain_core.documents", fromlist=["Document"]).Document(
                page_content="sheet",
                metadata={"page_number": 1},
            )
        ]
        docs = _load_legacy_office(str(path), "xls")

    assert len(docs) == 1
    mock_xlrd.assert_called_once_with(str(path))


def test_load_unsupported_type():
    with pytest.raises(ValueError, match="does not support"):
        load_with_python("/tmp/file.doc", "unknown")
